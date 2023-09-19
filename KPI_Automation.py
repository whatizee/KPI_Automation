# Description of the Porject : GSI IT OPS
# Purpose of the Script      : Daily KPI Report Generation
# Funtionality               : Fetch the Values from Excel and Paste it to PowerPoint
# Owner                      : GSI IT OPS
# Authur Bensil              : Y5S6ZV8 

import logging
import pptx
import pandas as pd
from pptx.util import Pt
from pptx import Presentation
from pptx.util import Inches
from termcolor import colored
from pptx.dml.color import RGBColor
import datetime
import time
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

yesterday = datetime.date.today() - datetime.timedelta(days=1)
formatted_date = yesterday.strftime("%d.%m.%Y")

print("Today's date in dd.mm.yyyy format:", formatted_date)
# start time
start_time = time.time()

def min_to_hr(minutes):
    hours = minutes / 60
    if hours.is_integer():
        return f"{hours:.0f}h"
    else:
        return f"{hours:.2f}h"
def min_to_min(minutes):
    return f"{minutes} mins"    
def convert_to_percent(percent):
    if percent.is_integer():
        return f"{percent:.0f}%"
    else:
        return f"{percent:.2f}%" 
def count_to_standard(number):
    if number > 999:
        return '{:,}'.format(number)
    else:
        return str(number)
        
sheet_names = ['New', 'New Monthly KPI', 'User feedback', 'New_ECP']
dfs = pd.read_excel('Daily Formulae_Feb.xlsx', sheet_name=sheet_names)

New_df = dfs['New']
New_df.fillna('', inplace=True)
New_Monthly_KPI_df = dfs['New Monthly KPI']
User_feedback_df = dfs['User feedback']
User_feedback_df.fillna('', inplace=True)
New_ECP = dfs['New_ECP']
New_ECP.fillna('', inplace=True)


#Fetching the last added value OVERALL RELOADED UPTIME
OR_Service_up_min_filtered_df = New_df[New_df.iloc[:, 16].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OR_Service_up_min_filtered_df.iloc[:, 16] = OR_Service_up_min_filtered_df.iloc[:, 16].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OR_Service_up_min_last_value = OR_Service_up_min_filtered_df.iloc[-1, 16]
OR_Service_up_hour_last_value = min_to_hr(OR_Service_up_min_last_value)

#Fetching the last added value OVERALL RELOADED DOWNTIME
OR_Service_down_min_filtered_df = New_df[New_df.iloc[:, 17].apply(lambda x: isinstance(x, (int, float)) and x >=0)]
OR_Service_down_min_filtered_df.iloc[:, 17] = OR_Service_down_min_filtered_df.iloc[:, 17].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OR_Service_down_min_last_value = OR_Service_down_min_filtered_df.iloc[-1, 17]
OR_Service_down_hour_last_value = min_to_min(OR_Service_down_min_last_value)

#Fetching the last added value Reloaded availability %
OR_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 1].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OR_Service_OA_percent_filtered_df.iloc[:, 1] = OR_Service_OA_percent_filtered_df.iloc[:, 1].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OR_Service_OA_percent_last_value = OR_Service_OA_percent_filtered_df.iloc[-1, 1]
OR_Service_OA_pecent_convert_last_value = convert_to_percent(OR_Service_OA_percent_last_value)

#Fetching the last added value OVERALL RELOADED Request calls
OR_Service_OA_call_filtered_df = New_df[New_df.iloc[:, 18].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OR_Service_OA_call_filtered_df.iloc[:, 18] = OR_Service_OA_call_filtered_df.iloc[:, 18].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OR_Service_OA_call_last_value = OR_Service_OA_call_filtered_df.iloc[-1, 18]
OR_Service_OA_call_standard_last_value = count_to_standard(OR_Service_OA_call_last_value)

#Fetching the last added value OVERALL RELOADED Request  failed calls
OR_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 19].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OR_Service_OA_call_failed_filtered_df.iloc[:, 19] = OR_Service_OA_call_failed_filtered_df.iloc[:, 19].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OR_Service_OA_call_failed_last_value = OR_Service_OA_call_failed_filtered_df.iloc[-1, 19]
OR_Service_OA_call_failed_standard_last_value = count_to_standard(OR_Service_OA_call_failed_last_value)

#Fetching the last added value Reloaded QB availability %
ORQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 2].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
ORQB_Service_OA_percent_filtered_df.iloc[:, 2] = ORQB_Service_OA_percent_filtered_df.iloc[:, 2].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORQB_Service_OA_percent_last_value = ORQB_Service_OA_percent_filtered_df.iloc[-1, 2]
ORQB_Service_OA_pecent_convert_last_value = convert_to_percent(ORQB_Service_OA_percent_last_value)


#WSI Overall
#Fetching the last added value OVERALL WSI UPTIME
OWSI_Service_up_min_filtered_df = New_df[New_df.iloc[:, 21].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSI_Service_up_min_filtered_df.iloc[:, 21] = OWSI_Service_up_min_filtered_df.iloc[:, 21].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSI_Service_up_min_last_value = OR_Service_up_min_filtered_df.iloc[-1, 21]
OWSI_Service_up_hour_last_value = min_to_hr(OWSI_Service_up_min_last_value)

#Fetching the last added value OVERALL WSI DOWNTIME
OWSI_Service_down_min_filtered_df = New_df[New_df.iloc[:, 22].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSI_Service_down_min_filtered_df.iloc[:, 22] = OWSI_Service_down_min_filtered_df.iloc[:, 22].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSI_Service_down_min_last_value = OWSI_Service_down_min_filtered_df.iloc[-1, 22]
OWSI_Service_down_hour_last_value = min_to_min(OWSI_Service_down_min_last_value)

#Fetching the last added value WSI availability %
OWSI_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 3].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OWSI_Service_OA_percent_filtered_df.iloc[:, 3] = OWSI_Service_OA_percent_filtered_df.iloc[:, 3].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSI_Service_OA_percent_last_value = OWSI_Service_OA_percent_filtered_df.iloc[-1, 3]
OWSI_Service_OA_pecent_convert_last_value = convert_to_percent(OWSI_Service_OA_percent_last_value)

#Fetching the last added value OVERALL WSI Request calls
OWSI_Service_OA_call_filtered_df = New_df[New_df.iloc[:, 23].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSI_Service_OA_call_filtered_df.iloc[:, 23] = OWSI_Service_OA_call_filtered_df.iloc[:, 23].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSI_Service_OA_call_last_value = OWSI_Service_OA_call_filtered_df.iloc[-1, 23]
OWSI_Service_OA_call_standard_last_value = count_to_standard(OWSI_Service_OA_call_last_value)

#Fetching the last added value OVERALL WSI Request  failed calls
OWSI_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 24].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OWSI_Service_OA_call_failed_filtered_df.iloc[:, 24] = OWSI_Service_OA_call_failed_filtered_df.iloc[:, 24].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSI_Service_OA_call_failed_last_value = OWSI_Service_OA_call_failed_filtered_df.iloc[-1, 24]
OWSI_Service_OA_call_failed_standard_last_value = count_to_standard(OWSI_Service_OA_call_failed_last_value)


#Fetching the last added value WSI request Success availability %
OWSIQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 4].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OWSIQB_Service_OA_percent_filtered_df.iloc[:, 4] = OWSIQB_Service_OA_percent_filtered_df.iloc[:, 4].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIQB_Service_OA_percent_last_value = OWSIQB_Service_OA_percent_filtered_df.iloc[-1, 4]
OWSIQB_Service_OA_pecent_convert_last_value = convert_to_percent(OWSIQB_Service_OA_percent_last_value)


#Reloaded individual tenents
#Fetching the last added value  FCA Q & B process RELOADED calls
ORFCAQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 34].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORFCAQB_Service_up_min_filtered_df.iloc[:, 34] = ORFCAQB_Service_up_min_filtered_df.iloc[:, 34].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFCAQB_Service_up_min_last_value = ORFCAQB_Service_up_min_filtered_df.iloc[-1, 34]
ORFCAQB_Service_up_hour_last_value = count_to_standard(ORFCAQB_Service_up_min_last_value)


#Fetching the last added value FCA Q & B process Request  failed calls
ORFCAQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 35].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORFCAQB_Service_OA_call_failed_filtered_df.iloc[:, 35] = ORFCAQB_Service_OA_call_failed_filtered_df.iloc[:, 35].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFCAQB_Service_OA_call_failed_last_value = ORFCAQB_Service_OA_call_failed_filtered_df.iloc[-1, 35]
ORFCAQB_Service_OA_call_failed_standard_last_value = count_to_standard(ORFCAQB_Service_OA_call_failed_last_value)


#Fetching the last added value F Q & B  Success request availability %
ORFCAQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 10].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
ORFCAQB_Service_OA_percent_filtered_df.iloc[:, 10] = ORFCAQB_Service_OA_percent_filtered_df.iloc[:, 10].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFCAQB_Service_OA_percent_last_value = ORFCAQB_Service_OA_percent_filtered_df.iloc[-1, 10]
ORFCAQB_Service_OA_pecent_convert_last_value = convert_to_percent(ORFCAQB_Service_OA_percent_last_value)

#Fetching the last added value  FORDQB process RELOADED calls
ORFORDQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 26].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORFORDQB_Service_up_min_filtered_df.iloc[:, 26] = ORFORDQB_Service_up_min_filtered_df.iloc[:, 26].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFORDQB_Service_up_min_last_value = ORFORDQB_Service_up_min_filtered_df.iloc[-1, 26]
ORFORDQB_Service_up_hour_last_value = count_to_standard(ORFORDQB_Service_up_min_last_value)

#Fetching the last added value FORDQB process Request  failed calls
ORFORDQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 27].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORFORDQB_Service_OA_call_failed_filtered_df.iloc[:, 27] = ORFORDQB_Service_OA_call_failed_filtered_df.iloc[:, 27].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFORDQB_Service_OA_call_failed_last_value = ORFORDQB_Service_OA_call_failed_filtered_df.iloc[-1, 27]
ORFORDQB_Service_OA_call_failed_standard_last_value = count_to_standard(ORFORDQB_Service_OA_call_failed_last_value)

#Fetching the last added value FORDQB  Success request availability %
ORFORDQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 6].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
ORFORDQB_Service_OA_percent_filtered_df.iloc[:, 6] = ORFORDQB_Service_OA_percent_filtered_df.iloc[:, 6].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORFORDQB_Service_OA_percent_last_value = ORFORDQB_Service_OA_percent_filtered_df.iloc[-1, 6]
ORFORDQB_Service_OA_pecent_convert_last_value = convert_to_percent(ORFORDQB_Service_OA_percent_last_value)

#Fetching the last added value  OVFQB process RELOADED calls
OROVFQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 32].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OROVFQB_Service_up_min_filtered_df.iloc[:, 32] = OROVFQB_Service_up_min_filtered_df.iloc[:, 32].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OROVFQB_Service_up_min_last_value = OROVFQB_Service_up_min_filtered_df.iloc[-1, 32]
OROVFQB_Service_up_hour_last_value = count_to_standard(OROVFQB_Service_up_min_last_value)

#Fetching the last added value OVFQB process Request  failed calls
OROVFQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 33].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OROVFQB_Service_OA_call_failed_filtered_df.iloc[:, 33] = OROVFQB_Service_OA_call_failed_filtered_df.iloc[:, 33].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OROVFQB_Service_OA_call_failed_last_value = OROVFQB_Service_OA_call_failed_filtered_df.iloc[-1, 33]
OROVFQB_Service_OA_call_failed_standard_last_value = count_to_standard(OROVFQB_Service_OA_call_failed_last_value)

#Fetching the last added value OVFQB  Success request availability %
OROVFQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 9].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OROVFQB_Service_OA_percent_filtered_df.iloc[:, 9] = OROVFQB_Service_OA_percent_filtered_df.iloc[:, 9].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OROVFQB_Service_OA_percent_last_value = OROVFQB_Service_OA_percent_filtered_df.iloc[-1, 9]
OROVFQB_Service_OA_pecent_convert_last_value = convert_to_percent(OROVFQB_Service_OA_percent_last_value)


#Fetching the last added value  PSAQB process RELOADED calls
ORPSAQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 30].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORPSAQB_Service_up_min_filtered_df.iloc[:, 30] = ORPSAQB_Service_up_min_filtered_df.iloc[:, 30].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORPSAQB_Service_up_min_last_value = ORPSAQB_Service_up_min_filtered_df.iloc[-1, 30]
ORPSAQB_Service_up_hour_last_value = count_to_standard(ORPSAQB_Service_up_min_last_value)

#Fetching the last added value PSAQB process Request  failed calls
ORPSAQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 31].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
ORPSAQB_Service_OA_call_failed_filtered_df.iloc[:, 31] = ORPSAQB_Service_OA_call_failed_filtered_df.iloc[:, 31].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORPSAQB_Service_OA_call_failed_last_value = ORPSAQB_Service_OA_call_failed_filtered_df.iloc[-1, 31]
ORPSAQB_Service_OA_call_failed_standard_last_value = count_to_standard(ORPSAQB_Service_OA_call_failed_last_value)

#Fetching the last added value PSAQB  Success request availability %
ORPSAQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 8].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
ORPSAQB_Service_OA_percent_filtered_df.iloc[:, 8] = ORPSAQB_Service_OA_percent_filtered_df.iloc[:, 8].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORPSAQB_Service_OA_percent_last_value = ORPSAQB_Service_OA_percent_filtered_df.iloc[-1, 8]
ORPSAQB_Service_OA_pecent_convert_last_value = convert_to_percent(ORPSAQB_Service_OA_percent_last_value)

#Fetching the last added value  PSAQB process RELOADED calls
ORVOQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 28].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
ORVOQB_Service_up_min_filtered_df.iloc[:, 28] = ORVOQB_Service_up_min_filtered_df.iloc[:, 28].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORVOQB_Service_up_min_last_value = ORVOQB_Service_up_min_filtered_df.iloc[-1, 28]
ORVOQB_Service_up_hour_last_value = count_to_standard(ORVOQB_Service_up_min_last_value)

#Fetching the last added value VOQB process Request  failed calls
ORVOQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 29].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
ORVOQB_Service_OA_call_failed_filtered_df.iloc[:, 29] = ORVOQB_Service_OA_call_failed_filtered_df.iloc[:, 29].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORVOQB_Service_OA_call_failed_last_value = ORVOQB_Service_OA_call_failed_filtered_df.iloc[-1, 29]
ORVOQB_Service_OA_call_failed_standard_last_value = count_to_standard(ORVOQB_Service_OA_call_failed_last_value)

#Fetching the last added value VOQB  Success request availability %
ORVOQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 7].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
ORVOQB_Service_OA_percent_filtered_df.iloc[:, 7] = ORVOQB_Service_OA_percent_filtered_df.iloc[:, 7].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
ORVOQB_Service_OA_percent_last_value = ORVOQB_Service_OA_percent_filtered_df.iloc[-1, 7]
ORVOQB_Service_OA_pecent_convert_last_value = convert_to_percent(ORVOQB_Service_OA_percent_last_value)


#Fetching the last added value  AVSQB process WSI calls
OWSIAVSQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 39].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSIAVSQB_Service_up_min_filtered_df.iloc[:, 39] = OWSIAVSQB_Service_up_min_filtered_df.iloc[:, 39].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIAVSQB_Service_up_min_last_value = OWSIAVSQB_Service_up_min_filtered_df.iloc[-1, 39]
OWSIAVSQB_Service_up_hour_last_value = count_to_standard(OWSIAVSQB_Service_up_min_last_value)

#Fetching the last added value WSIAVSQB process Request  failed calls
OWSIAVSQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 40].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OWSIAVSQB_Service_OA_call_failed_filtered_df.iloc[:, 40] = OWSIAVSQB_Service_OA_call_failed_filtered_df.iloc[:, 40].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIAVSQB_Service_OA_call_failed_last_value = OWSIAVSQB_Service_OA_call_failed_filtered_df.iloc[-1, 40]
OWSIAVSQB_Service_OA_call_failed_standard_last_value = count_to_standard(OWSIAVSQB_Service_OA_call_failed_last_value)

#Fetching the last added value WSIAVSQB  Success request availability %
OWSIAVSQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 13].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OWSIAVSQB_Service_OA_percent_filtered_df.iloc[:, 13] = OWSIAVSQB_Service_OA_percent_filtered_df.iloc[:, 13].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIAVSQB_Service_OA_percent_last_value = OWSIAVSQB_Service_OA_percent_filtered_df.iloc[-1, 13]
OWSIAVSQB_Service_OA_pecent_convert_last_value = convert_to_percent(OWSIAVSQB_Service_OA_percent_last_value)


#Fetching the last added value  WSIVMQB process WSI calls
OWSIVMQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 37].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSIVMQB_Service_up_min_filtered_df.iloc[:, 37] = OWSIVMQB_Service_up_min_filtered_df.iloc[:, 37].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIVMQB_Service_up_min_last_value = OWSIVMQB_Service_up_min_filtered_df.iloc[-1, 37]
OWSIVMQB_Service_up_hour_last_value = count_to_standard(OWSIVMQB_Service_up_min_last_value)

#Fetching the last added value WSIVMQB process Request  failed calls
OWSIVMQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 38].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OWSIVMQB_Service_OA_call_failed_filtered_df.iloc[:, 38] = OWSIVMQB_Service_OA_call_failed_filtered_df.iloc[:, 38].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIVMQB_Service_OA_call_failed_last_value = OWSIVMQB_Service_OA_call_failed_filtered_df.iloc[-1, 38]
OWSIVMQB_Service_OA_call_failed_standard_last_value = count_to_standard(OWSIVMQB_Service_OA_call_failed_last_value)

#Fetching the last added value WSIVMQB  Success request availability %
OWSIVMQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 12].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OWSIVMQB_Service_OA_percent_filtered_df.iloc[:, 12] = OWSIVMQB_Service_OA_percent_filtered_df.iloc[:, 12].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIVMQB_Service_OA_percent_last_value = OWSIVMQB_Service_OA_percent_filtered_df.iloc[-1, 12]
OWSIVMQB_Service_OA_pecent_convert_last_value = convert_to_percent(OWSIVMQB_Service_OA_percent_last_value)


#Fetching the last added value  WSIOVFQB process WSI calls
OWSIOVFQB_Service_up_min_filtered_df = New_df[New_df.iloc[:, 41].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OWSIOVFQB_Service_up_min_filtered_df.iloc[:, 41] = OWSIOVFQB_Service_up_min_filtered_df.iloc[:, 41].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIOVFQB_Service_up_min_last_value = OWSIOVFQB_Service_up_min_filtered_df.iloc[-1, 41]
OWSIOVFQB_Service_up_hour_last_value = count_to_standard(OWSIOVFQB_Service_up_min_last_value)

#Fetching the last added value WSIOVFQB process Request  failed calls
OWSIOVFQB_Service_OA_call_failed_filtered_df = New_df[New_df.iloc[:, 42].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OWSIOVFQB_Service_OA_call_failed_filtered_df.iloc[:, 42] = OWSIOVFQB_Service_OA_call_failed_filtered_df.iloc[:, 42].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIOVFQB_Service_OA_call_failed_last_value = OWSIOVFQB_Service_OA_call_failed_filtered_df.iloc[-1, 42]
OWSIOVFQB_Service_OA_call_failed_standard_last_value = count_to_standard(OWSIOVFQB_Service_OA_call_failed_last_value)

#Fetching the last added value WSIOVFQB  Success request availability %
OWSIOVFQB_Service_OA_percent_filtered_df = New_df[New_df.iloc[:, 14].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OWSIOVFQB_Service_OA_percent_filtered_df.iloc[:, 14] = OWSIOVFQB_Service_OA_percent_filtered_df.iloc[:, 14].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OWSIOVFQB_Service_OA_percent_last_value = OWSIOVFQB_Service_OA_percent_filtered_df.iloc[-1, 14]
OWSIOVFQB_Service_OA_pecent_convert_last_value = convert_to_percent(OWSIOVFQB_Service_OA_percent_last_value)



#Fetching the last added value OVERALL ECP UPTIME
OECP_Service_up_min_filtered_df = New_ECP[New_ECP.iloc[:, 14].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OECP_Service_up_min_filtered_df.iloc[:, 14] = OECP_Service_up_min_filtered_df.iloc[:, 14].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECP_Service_up_min_last_value = OECP_Service_up_min_filtered_df.iloc[-1, 14]
OECP_Service_up_hour_last_value = min_to_hr(OECP_Service_up_min_last_value)

#Fetching the last added value ECP  DOWNTIME
OECP_Service_down_min_filtered_df = New_ECP[New_ECP.iloc[:, 15].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OECP_Service_down_min_filtered_df.iloc[:, 15] = OECP_Service_down_min_filtered_df.iloc[:, 15].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECP_Service_down_min_last_value = OECP_Service_down_min_filtered_df.iloc[-1, 15]
OECP_Service_down_hour_last_value = min_to_min(OECP_Service_down_min_last_value)

#Fetching the last added value ECP availability %
OECP_Service_OA_percent_filtered_df = New_ECP[New_ECP.iloc[:, 1].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OECP_Service_OA_percent_filtered_df.iloc[:, 1] = OECP_Service_OA_percent_filtered_df.iloc[:, 1].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECP_Service_OA_percent_last_value = OECP_Service_OA_percent_filtered_df.iloc[-1, 1]
OECP_Service_OA_pecent_convert_last_value = convert_to_percent(OECP_Service_OA_percent_last_value)


#Fetching the last added value  ECPQB process WSI calls
OECPQB_Service_up_min_filtered_df = New_ECP[New_ECP.iloc[:, 16].apply(lambda x: isinstance(x, (int, float)) and x >= 0)]
OECPQB_Service_up_min_filtered_df.iloc[:, 16] = OECPQB_Service_up_min_filtered_df.iloc[:, 16].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECPQB_Service_up_min_last_value = OECPQB_Service_up_min_filtered_df.iloc[-1, 16]
OECPQB_Service_up_hour_last_value = count_to_standard(OECPQB_Service_up_min_last_value)

#Fetching the last added value ECPQB process Request  failed calls
OECPQB_Service_OA_call_failed_filtered_df = New_ECP[New_ECP.iloc[:, 17].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OECPQB_Service_OA_call_failed_filtered_df.iloc[:, 17] = OECPQB_Service_OA_call_failed_filtered_df.iloc[:, 17].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECPQB_Service_OA_call_failed_last_value = OECPQB_Service_OA_call_failed_filtered_df.iloc[-1, 17]
OECPQB_Service_OA_call_failed_standard_last_value = count_to_standard(OECPQB_Service_OA_call_failed_last_value)

#Fetching the last added value ECPQB  Success request availability %
OECPQB_Service_OA_percent_filtered_df = New_ECP[New_ECP.iloc[:, 2].apply(lambda x: isinstance(x, (int, float)) and x > 0)]
OECPQB_Service_OA_percent_filtered_df.iloc[:, 2] = OECPQB_Service_OA_percent_filtered_df.iloc[:, 2].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OECPQB_Service_OA_percent_last_value = OECPQB_Service_OA_percent_filtered_df.iloc[-1, 2]
OECPQB_Service_OA_pecent_convert_last_value = convert_to_percent(OECPQB_Service_OA_percent_last_value)

#Removing last two rows of dataframe
User_feedback_df = User_feedback_df.iloc[:-2]

#fetching the FCA USER feedback total number 
FCAUFB_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 1].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
FCAUFB_Service_up_min_filtered_df.iloc[:, 1] = FCAUFB_Service_up_min_filtered_df.iloc[:, 1].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
FCAUFB_Service_up_min_last_value = FCAUFB_Service_up_min_filtered_df.iloc[-1, 1]
FCAUFB_Service_OA_feedback_standard_last_value = count_to_standard(FCAUFB_Service_up_min_last_value)
print(FCAUFB_Service_OA_feedback_standard_last_value)

#fetching the FCA USER feedback total Rating 
FCAUFR_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 2].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
FCAUFR_Service_up_min_filtered_df.iloc[:, 2] = FCAUFR_Service_up_min_filtered_df.iloc[:, 2].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
FCAUFR_Service_up_min_last_value = FCAUFR_Service_up_min_filtered_df.iloc[-1, 2]
FCAUFR_Service_OA_feedback_standard_last_value = count_to_standard(FCAUFR_Service_up_min_last_value)

#fetching the FORD USER feedback total number 
FORDUFB_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 3].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
FORDUFB_Service_up_min_filtered_df.iloc[:, 3] = FORDUFB_Service_up_min_filtered_df.iloc[:, 3].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
FORDUFB_Service_up_min_last_value = FCAUFB_Service_up_min_filtered_df.iloc[-1, 3]
FORDUFB_Service_OA_feedback_standard_last_value = count_to_standard(FORDUFB_Service_up_min_last_value)


#fetching the FORDUFR USER feedback total Rating 
FORDUFR_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 4].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
FORDUFR_Service_up_min_filtered_df.iloc[:, 4] = FORDUFR_Service_up_min_filtered_df.iloc[:, 4].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
FORDUFR_Service_up_min_last_value = FCAUFR_Service_up_min_filtered_df.iloc[-1, 4]
FORDUFR_Service_OA_feedback_standard_last_value = count_to_standard(FORDUFR_Service_up_min_last_value)

#fetching the OVF USER feedback total number 
OVFUFB_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 5].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OVFUFB_Service_up_min_filtered_df.iloc[:, 5] = OVFUFB_Service_up_min_filtered_df.iloc[:, 5].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OVFUFB_Service_up_min_last_value = OVFUFB_Service_up_min_filtered_df.iloc[-1, 5]
OVFUFB_Service_OA_feedback_standard_last_value = count_to_standard(OVFUFB_Service_up_min_last_value)


#fetching the OVF USER feedback total Rating 
OVFUFR_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 6].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
OVFUFR_Service_up_min_filtered_df.iloc[:, 6] = OVFUFR_Service_up_min_filtered_df.iloc[:, 6].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
OVFUFR_Service_up_min_last_value = OVFUFR_Service_up_min_filtered_df.iloc[-1, 6]
OVFUFR_Service_OA_feedback_standard_last_value = count_to_standard(OVFUFR_Service_up_min_last_value)


#fetching the PSAUFB USER feedback total number 
PSAUFB_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 7].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
PSAUFB_Service_up_min_filtered_df.iloc[:, 7] = PSAUFB_Service_up_min_filtered_df.iloc[:, 7].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
PSAUFB_Service_up_min_last_value = PSAUFB_Service_up_min_filtered_df.iloc[-1, 7]
PSAUFB_Service_OA_feedback_standard_last_value = count_to_standard(PSAUFB_Service_up_min_last_value)


#fetching the PSAUFR USER feedback total Rating 
PSAUFR_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 8].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
PSAUFR_Service_up_min_filtered_df.iloc[:, 8] = PSAUFR_Service_up_min_filtered_df.iloc[:, 8].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
PSAUFR_Service_up_min_last_value = PSAUFR_Service_up_min_filtered_df.iloc[-1, 8]
PSAUFR_Service_OA_feedback_standard_last_value = count_to_standard(PSAUFR_Service_up_min_last_value)


#fetching the VOLUFB USER feedback total number 
VOLUFB_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 9].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
VOLUFB_Service_up_min_filtered_df.iloc[:, 9] = VOLUFB_Service_up_min_filtered_df.iloc[:, 9].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
VOLUFB_Service_up_min_last_value = VOLUFB_Service_up_min_filtered_df.iloc[-1, 9]
VOLUFB_Service_OA_feedback_standard_last_value = count_to_standard(VOLUFB_Service_up_min_last_value)


#fetching the VOLUFR USER feedback total Rating 
VOLUFR_Service_up_min_filtered_df = User_feedback_df[User_feedback_df.iloc[:, 10].apply(lambda x: isinstance(x, (int, float))and x >= 0)]
VOLUFR_Service_up_min_filtered_df.iloc[:, 10] = VOLUFR_Service_up_min_filtered_df.iloc[:, 10].apply(lambda x: x.replace('\n', '') if isinstance(x, str) else x)
VOLUFR_Service_up_min_last_value = VOLUFR_Service_up_min_filtered_df.iloc[-1, 10]
VOLUFR_Service_OA_feedback_standard_last_value = count_to_standard(VOLUFR_Service_up_min_last_value)

# Load the PowerPoint file
ppt_file = pptx.Presentation('D:\KPI_Automation\CrossCheck\Daily Portal KPI Report.pptx')



# Get the first slide
slide = ppt_file.slides[0]

# Get the shape you want to replace
shape = slide.shapes[0]

# Adding Value to date
#Dimensions of the text box
left = Inches(10.5)
top = Inches(.350)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = 'Status ' + formatted_date 
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(11)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.paragraphs[0].font.bold = True


# Adding Value to OVERALL RELOADED UPTIME
OR_Service_up_hour_last_string_value = str(OR_Service_up_hour_last_value)
#Dimensions of the text box
left = Inches(2.2)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OR_Service_up_hour_last_string_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL RELOADED Downtime
OR_Service_down_hour_last_string_value = str(OR_Service_down_hour_last_value)
#Dimensions of the text box
left = Inches(2.1)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OR_Service_down_hour_last_string_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL RELOADED percentage
#Dimensions of the text box
left = Inches(2.057)
top = Inches(2.57)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OR_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
# Set the text box properties
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OR_Service_OA_percent_last_value > 99.80 and OR_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OR_Service_OA_percent_last_value > 97.00 and OR_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color
    
# Adding Value to OVERALL RELOADED request calls   
#Dimensions of the text box
left = Inches(4.1)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OR_Service_OA_call_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL RELOADED Failed calls
#Dimensions of the text box
left = Inches(4.2)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OR_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL RELOADED Request percentage
#Dimensions of the text box
left = Inches(4.04)
top = Inches(2.56)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if ORQB_Service_OA_percent_last_value > 99.80 and ORQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif ORQB_Service_OA_percent_last_value > 97.00 and ORQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color
    
# Adding Value to OVERALL WSI UPTIME
OWSI_Service_up_hour_last_string_value = str(OWSI_Service_up_hour_last_value)
#Dimensions of the text box
left = Inches(8)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSI_Service_up_hour_last_string_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL WSI Downtime
#Dimensions of the text box
left = Inches(8)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSI_Service_down_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL WSI percentage
#Dimensions of the text box
left = Inches(7.97)
top = Inches(2.56)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSI_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OWSI_Service_OA_percent_last_value > 99.80 and OWSI_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OWSI_Service_OA_percent_last_value > 97.00 and OWSI_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color

# Adding Value to OVERALL WSI request calls   
#Dimensions of the text box
left = Inches(10)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSI_Service_OA_call_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL WSI Failed calls
#Dimensions of the text box
left = Inches(10.1)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSI_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to OVERALL WSI Request percentage
#Dimensions of the text box
left = Inches(9.92)
top = Inches(2.56)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OWSIQB_Service_OA_percent_last_value > 99.80 and OWSIQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OWSIQB_Service_OA_percent_last_value > 97.00 and OWSIQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color
    
# Adding Value to  FCAQB Reloaed request calls   
#Dimensions of the text box
left = Inches(2.2)
top = Inches(3.6)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFCAQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to FCAQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(2.2)
top = Inches(4.0)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFCAQB_Service_OA_call_failed_standard_last_value
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to FCAQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(2.05)
top = Inches(4.34)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFCAQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if ORFCAQB_Service_OA_percent_last_value > 99.80 and ORFCAQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif ORFCAQB_Service_OA_percent_last_value > 97.00 and ORFCAQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color
    
    
# Adding Value to  FORDQB Reloaed request calls   
#Dimensions of the text box
left = Inches(4.2)
top = Inches(3.6)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFORDQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to FORDQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(4.2)
top = Inches(4.0)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFORDQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to FORDQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(4.03)
top = Inches(4.34)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORFORDQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if ORFORDQB_Service_OA_percent_last_value > 99.80 and ORFORDQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif ORFORDQB_Service_OA_percent_last_value > 97.00 and ORFORDQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color


# Adding Value to  OVFQB Reloaed request calls   
#Dimensions of the text box
left = Inches(6.2)
top = Inches(3.6)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OROVFQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVFQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(6.2)
top = Inches(4.0)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OROVFQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to OVFQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(6)
top = Inches(4.34)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OROVFQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OROVFQB_Service_OA_percent_last_value > 99.80 and OROVFQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OROVFQB_Service_OA_percent_last_value > 97.00 and OROVFQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    

# Adding Value to  PSAQB Reloaed request calls   
#Dimensions of the text box
left = Inches(8.2)
top = Inches(3.6)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORPSAQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to PSAQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(8.2)
top = Inches(4.0)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORPSAQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to PSAQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(7.95)
top = Inches(4.34)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORPSAQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if ORPSAQB_Service_OA_percent_last_value > 99.80 and ORPSAQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif ORPSAQB_Service_OA_percent_last_value > 97.00 and ORPSAQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color 

    
# Adding Value to  VOQB Reloaed request calls   
#Dimensions of the text box
left = Inches(10)
top = Inches(3.6)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORVOQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to VOQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(10)
top = Inches(4.0)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORVOQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to VOQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(9.92)
top = Inches(4.34)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = ORVOQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if ORVOQB_Service_OA_percent_last_value > 99.80 and ORVOQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif ORVOQB_Service_OA_percent_last_value > 97.00 and ORVOQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    


# Adding Value to  WSIAVSQB Reloaed request calls   
#Dimensions of the text box
left = Inches(2.1)
top = Inches(5.4)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIAVSQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to WSIAVSQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(2.2)
top = Inches(5.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIAVSQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to WSIAVSQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(2.05)
top = Inches(6.17)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIAVSQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OWSIAVSQB_Service_OA_percent_last_value > 99.80 and OWSIAVSQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OWSIAVSQB_Service_OA_percent_last_value > 97.00 and OWSIAVSQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    


# Adding Value to  WSIVMQB Reloaed request calls   
#Dimensions of the text box
left = Inches(4.1)
top = Inches(5.4)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIVMQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to WSIVMQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(4.2)
top = Inches(5.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIVMQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to WSIVMQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(4.02)
top = Inches(6.17)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIVMQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OWSIVMQB_Service_OA_percent_last_value > 99.80 and OWSIVMQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OWSIVMQB_Service_OA_percent_last_value > 97.00 and OWSIVMQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    

# Adding Value to  WSIOVFQB Reloaed request calls   
#Dimensions of the text box
left = Inches(6.1)
top = Inches(5.4)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIOVFQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to WSIOVFQB RELOADED Failed calls
#Dimensions of the text box
left = Inches(6.2)
top = Inches(5.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIOVFQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)    


# Adding Value to WSIOVFQB RELOADED Request percentage
#Dimensions of the text box
left = Inches(5.98)
top = Inches(6.18)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OWSIOVFQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OWSIOVFQB_Service_OA_percent_last_value > 99.80 and OWSIOVFQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OWSIOVFQB_Service_OA_percent_last_value > 97.00 and OWSIOVFQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    

# Get the second slide
slide = ppt_file.slides[1]

# Get the shape you want to replace
shape = slide.shapes[0]

# Adding Value to date
#Dimensions of the text box
left = Inches(10.5)
top = Inches(.350)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = 'Status ' + formatted_date 
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(11)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.paragraphs[0].font.bold = True    
    
# Adding Value to OVERALL ECP UPTIME
#Dimensions of the text box
left = Inches(2.2)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECP_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL ECP Downtime
#Dimensions of the text box
left = Inches(2.1)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECP_Service_down_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL ECP percentage
#Dimensions of the text box
left = Inches(2.05)
top = Inches(2.57)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECP_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OECP_Service_OA_percent_last_value > 99.80 and OECP_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OECP_Service_OA_percent_last_value > 97.00 and OECP_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    

# Adding Value to OVERALL ECP call success
#Dimensions of the text box
left = Inches(4.2)
top = Inches(1.75)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to OVERALL ECP request failer Downtime
#Dimensions of the text box
left = Inches(4.1)
top = Inches(2.275)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to  Overall ECP  request percentage
#Dimensions of the text box
left = Inches(4.05)
top = Inches(2.57)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OECPQB_Service_OA_percent_last_value > 99.80 and OECPQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OECPQB_Service_OA_percent_last_value > 97.00 and OECPQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color    

# Adding Value to  ECP requests
#Dimensions of the text box
left = Inches(2.2)
top = Inches(3.5)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_up_hour_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to  ECP request failure
#Dimensions of the text box
left = Inches(2.1)
top = Inches(4)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_OA_call_failed_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)

# Adding Value to  ECP request sucess percentage
#Dimensions of the text box
left = Inches(2.06)
top = Inches(4.35)
width = Inches(.68)
height = Inches(.4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OECPQB_Service_OA_pecent_convert_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
if OECPQB_Service_OA_percent_last_value > 99.80 and OECPQB_Service_OA_percent_last_value <= 100:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(144, 238, 144)  # green color
elif OECPQB_Service_OA_percent_last_value > 97.00 and OECPQB_Service_OA_percent_last_value <= 99.79:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 191)  # amber color
else:
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 117, 117)  # red color   


# Get the second slide
slide = ppt_file.slides[2]

# Get the shape you want to replace
shape = slide.shapes[0]

# Adding Value to date
#Dimensions of the text box
left = Inches(10.5)
top = Inches(.350)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = 'Status ' + formatted_date 
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(11)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)
text_box.text_frame.paragraphs[0].font.bold = True    

# Adding Value to  FCAQB User feedback
#Dimensions of the text box
left = Inches(1.2)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = FCAUFB_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   


# Adding Value to  FCAQB User feedback
#Dimensions of the text box
left = Inches(2.2)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = FCAUFR_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   

# Adding Value to  FORDUFB User feedback
#Dimensions of the text box
left = Inches(3.3)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = FORDUFB_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   


# Adding Value to  FORDUFR User Rating
#Dimensions of the text box
left = Inches(4.3)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = FORDUFR_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   

# Adding Value to  OVFQB User feedback
#Dimensions of the text box
left = Inches(5.45)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OVFUFB_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   


# Adding Value to  OVFQR User Rating
#Dimensions of the text box
left = Inches(6.45)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = OVFUFR_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   

# Adding Value to  PSAUFB User feedback
#Dimensions of the text box
left = Inches(7.56)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = PSAUFB_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   


# Adding Value to  PSAUFR User Rating
#Dimensions of the text box
left = Inches(8.56)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = PSAUFR_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139) 


# Adding Value to  VOLUFB User feedback
#Dimensions of the text box
left = Inches(9.7)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = VOLUFB_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)   


# Adding Value to  VOLUFR User Rating
#Dimensions of the text box
left = Inches(10.7)
top = Inches(2.8)
width = Inches(.3)
height = Inches(.3)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.text = VOLUFR_Service_OA_feedback_standard_last_value
# Set the font size of the text box based on the value in the second column of the row
text_box.text_frame.paragraphs[0].font.size = Pt(10)
text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139) 
       
# Save the Daily Portal KPI Report PowerPoint file
ppt_file.save('Daily Portal KPI Report' + '-'+ formatted_date + '.pptx')

# end time
end_time = time.time()

# calculate the execution time
execution_time = end_time - start_time

print("***********************************************************************")
print("Script Successfully excecuted")
# print the execution time
print(f"Script execution time: {execution_time:.2f} seconds")
print("***********************************************************************")